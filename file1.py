# интересные программы с курса основы программирования на языке python
import xlsxwriter
import sys

workbook = xlsxwriter.Workbook("res.xlsx")
worksheet = workbook.add_worksheet()  # add_worksheet надо добавить страницу в документ.
data = [i.split()(i[0], int(i[1])) for i in sys.stdin.read().split("\n")[:-2]]
for row, (item, price) in enumerate(data):
    worksheet.write(
        row, 0, item
    )  # write, который записывает в опред-ю строку и колонку данные.
    worksheet.write(row, 1, price)
chart = workbook.add_chart({"type": "pie"})
chart.add_series({"values": "=Sheet1!$B$1:$B$5"})
worksheet.insert_chart("C3", chart)
workbook.close()

# -----------------------------------------------------
from pptx import Presentation
import random

# создаем новую презентацию
prs = Presentation("res.pptx")

for method_name in dir(random):
    if not method_name.startswith("_") and method_name != "BPF":
        title_s = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_s)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = f"Метод {method_name}"
        subtitle.text = "\n".join(help(getattr(random, method_name)).split("\n")[2:])
prs.save("res.pptx")
# -------------------------------------------
import pymorphy2

morph = pymorphy2.MorphAnalyzer()
data = input()
if morph.parse(data)[0].tag.POS in {"INFN", "VERB"}:
    a = morph.parse(data)[0]
    print("Прошедшее время:")
    for i in ["masc", "femn", "neut", "plur"]:
        print(a.inflect({"past", i}).word)
    print("Настоящее время:")
    for i in ["1per", "2per", "3per"]:
        for j in ["sing", "plur"]:
            print(a.inflect({j, "pres", i}).word)
else:
    print("Не глагол")

# -----------------------------
import numpy as np
from PIL import Image


def bw_convert():
    img = Image.open('image.jpg')
    arr = np.asarray(img, dtype='uint8')
    x, y, _ = arr.shape
    k = np.array([[[0.2989, 0.587, 0.114]]])
    sums = np.round(np.sum(arr * k, axis=2)).astype(np.uint8)
    arr2 = np.repeat(sums[:, :, np.newaxis], 3, axis=2)
    img2 = Image.fromarray(arr2)
    img2.save('res.jpg')



# -------------------------------
import wave
import struct


def pitch_and_toss():
    source = wave.open("in.wav", mode="rb")
    dest = wave.open("out.wav", mode="wb")

    dest.setparams(source.getparams())
    # количество фреймов(1 амплитуда)
    frames_count = source.getnframes()
    data = struct.unpack("<" + str(frames_count) + "h", source.readframes(frames_count))

    newdata = frames_count // 4

    a, s, d, f = (
        data[:newdata],
        data[newdata: (2 * newdata)],
        data[2 * newdata: (3 * newdata)],
        data[3 * newdata:],
    )

    answer = d + f + a + s
    answer = struct.pack("<" + str(len(data)) + "h", *answer)
    dest.writeframes(answer)
    source.close()
    dest.close()



import wave
import struct


def chip_and_dale(number):
    source = wave.open("in.wav", mode="rb")
    dest = wave.open("out.wav", mode="wb")

    dest.setparams(source.getparams())
    # найдем количество фреймов
    frames_count = source.getnframes()
    data = struct.unpack("<" + str(frames_count) + "h", source.readframes(frames_count))

    # собственно, основная строка программы - переворот списка
    newdata = data[::number]

    newframes = struct.pack("<" + str(len(newdata)) + "h", *newdata)

    # записываем содержимое в преобразованный файл.
    dest.writeframes(newframes)
    source.close()
    dest.close()


chip_and_dale(3)

# -----------------------------
def transparency(filename1, filename2):
    from PIL import Image

    im_1 = Image.open(filename1)
    im_2 = Image.open(filename2)
    pixels1 = im_1.load()
    pixels2 = im_2.load()
    x, y = im_2.size
    for i in range(x):
        for j in range(y):
            r1, g1, b1 = pixels1[i, j]
            r2, g2, b2 = pixels2[i, j]
            r = int(0.5 * r1 + 0.5 * r2)
            g = int(0.5 * g1 + 0.5 * g2)
            b = int(0.5 * b1 + 0.5 * b2)
            pixels2[i, j] = r, g, b
    im_2.save('res.jpg')
    im_2.show()


# ----------------------

from PIL import Image


def makeanagliph(filename, delta):
    im = Image.open(filename)
    x, y = im.size
    res = Image.new('RGB', (x, y), (0, 0, 0))
    pixels2 = res.load()
    pixels = [(i % x, i // x) for i in range(x * y)]
    for i, j in pixels:
        r, g, b = im.getpixel((i, j))
        if i < delta:
            pixels2[i, j] = (0, g, b)
        else:
            r, _, _ = im.getpixel((i - delta, j))
            pixels2[i, j] = (r, g, b)
    res.save("res.jpg")
# -------------------------------------
from PIL import Image, ImageDraw

# Создание нового изображения
image = Image.new("RGB", (500, 200), (255, 255, 255))

# Создание объекта для рисования
draw = ImageDraw.Draw(image)


def gradient(color):
    new_color = (0, 0, 0)
    new_image = Image.new("RGB", (712, 200), new_color)
    draw = ImageDraw.Draw(new_image)
    pixels = new_image.load()
    r, g, b = 0, 0, 0
    color = color.lower()
    if color == 'r':
        for i in range(512):
            draw.line((i, 0, i, 200), fill=(i // 2, 0, 0), width=1)
            # меняем обциссу
    elif color == 'g':
        for j in range(512):
            draw.line((j, 0, j, 200), fill=(0, j // 2, 0), width=1)
    elif color == 'b':
        for k in range(512):
            draw.line((k, 0, k, 200), fill=(0, 0, k // 2), width=1)
    
    # Нарисовать букву 'А'
    draw.line((50, 50, 50, 150), fill=(255, 0, 0), width=10)
    draw.line((50, 50, 100, 50), fill=(255, 0, 0), width=10)
    draw.line((100, 50, 100, 150), fill=(255, 0, 0), width=10)
    draw.line((50, 100, 100, 100), fill=(255, 0, 0), width=10)

    draw.line((100, 50, 100, 150), fill=(255, 255, 255), width=10)
    draw.line((50, 100, 100, 100), fill=(255, 255, 255), width=10)

    # Нарисовать букву 'Л'
    draw.line((175, 50, 175, 150), fill=(0, 255, 0), width=10)
    draw.line((175, 150, 225, 150), fill=(0, 255, 0), width=10)

    # Нарисовать букву 'н'
    draw.line((300, 50, 300, 150), fill=(255, 255, 255), width=10)
    
    draw.line((425, 50, 425, 150), fill=(255, 255, 255), width=10)
    draw.line((475, 50, 475, 150), fill=(255, 255, 0), width=10)
    draw.line((425, 100, 475, 100), fill=(255, 255, 0), width=10)

    draw.line((600, 50, 600, 150), fill=(255, 255, 255), width=10)
    draw.line((600, 50, 650, 50), fill=(255, 255, 255), width=10)
    draw.line((650, 50, 650, 150), fill=(255, 255, 255), width=10)
    draw.line((600, 100, 650, 100), fill=(255, 255, 255), width=10)
    
    new_image.save("res.png")
    new_image.show()
    return


gradient('R')
# Сохранить изображение в файл
image.save("alya.png")

# ----
from PIL import Image, ImageDraw
 
 
def gradient(color):
    new_color = (0, 0, 0)
    new_image = Image.new("RGB", (512, 200), new_color)
    draw = ImageDraw.Draw(new_image)
    pixels = new_image.load()
    r, g, b = 0, 0, 0
    color = color.lower()
    if color == 'r':
        for i in range(512):
            draw.line((i, 0, i, 200), fill=(i // 2, 0, 0), width=1)
            # меняем обциссу
    elif color == 'g':
        for j in range(512):
            draw.line((j, 0, j, 200), fill=(0, j // 2, 0), width=1)
    elif color == 'b':
        for k in range(512):
            draw.line((k, 0, k, 200), fill=(0, 0, k // 2), width=1)
    
    new_image.save("res.png")
    new_image.show()
    return


gradient('R')


# ---------------------------
def cached(fun):
    dct = {}

    def inside(*arg):
        nonlocal dct
        # ускоорение из-за того, что если н есть в словаре,
        # то повторно иметь не будем
        if arg in dct:
            return dct[arg]
        # если нет, то запишем -это в словарь, но движуха в основном cashed
        # значит вызывем функцию fun в словарь, который nоnlkcal
        magic = fun(*arg)
        dct[arg] = magic
        return magic
    return inside


@cached
def fib(n):
    if n == 1 or n == 2:
        return 1
    else:
        return fib(n - 1) + fib(n - 2)
    
# ------------------------------------------
import tkinter
import random
from turtle import *

def move_wrap(obj, move):
    canvas.move(obj, move[0], move[1])
    x = canvas.coords(obj)[0]
    y = canvas.coords(obj)[1]
    if x < 0:
        canvas.move(obj, 600, 0)
    if x > 600:
        canvas.move(obj, -600, 0)
    if y > 600:
        canvas.move(obj, 0, -600)
    if y < 0:
        canvas.move(obj, 0, 600)
    

def do_nothing(x):
    pass

#------------------------------------------------
def check_move():
    if canvas.coords(player) == canvas.coords(exit):
        label.config(text="Победа!")
        master.bind("<KeyPress>", do_nothing)
        heart()
    for f in fires:
        if canvas.coords(player) == canvas.coords(f):
            label.config(text="Ты проиграл!")
            master.bind("<KeyPress>", do_nothing)


def heart():
    color ("pink")
    begin_fill()
    pensize(3)
    left(50)
    forward(133)
    circle(50,200)
    right(140)
    circle (50,200)
    forward (133)
    end_fill ()


def prepare_and_start():
    global player, exit, fires
    canvas.delete("all")
    player_pos = (random.randint(0, N_X - 1) * step,
                  random.randint(0, N_Y - 1) * step)
    exit_pos = (random.randint(0, N_X - 1) * step, 
                random.randint(0, N_Y - 1) * step)
    exit = canvas.create_oval(
                            (exit_pos[0], exit_pos[1]),
                            (exit_pos[0] + step, exit_pos[1] + step),
                            fill='yellow'
                            )
    player = canvas.create_oval(
                                (player_pos[0], player_pos[1]),
                                (player_pos[0] + step, player_pos[1] + step),
                                fill='green'
                                )
    N_FIRES = 6 
    fires = []
    for i in range(N_FIRES):
        fire_pos = (random.randint(0, N_X - 1) * step, 
                    random.randint(0, N_Y - 1) * step)
        fire = canvas.create_oval(
            (fire_pos[0], fire_pos[1]),
            (fire_pos[0] + step, fire_pos[1] + step),
            fill='red')
        fires.append(fire)

    label.config(text="Найди выход!")
    master.bind("<KeyPress>", key_pressed)


def key_pressed(event):
    if event.keysym == 'Up':
        move_wrap(player, (0, -step))
    elif event.keysym == 'Left':    move_wrap(player, (-step, 0))
    elif event.keysym == 'Right':   move_wrap(player, (step, 0))
    elif event.keysym == 'Down':   move_wrap(player, (0, step))
    check_move()

master = tkinter.Tk()
step = 60 
N_X = 10
N_Y = 10 
master = tkinter.Tk()
label = tkinter.Label(master, text="Найди выход")
label.pack()
canvas = tkinter.Canvas(master, bg='blue', 
                        height=N_X * step, width=N_Y * step)
canvas.pack()
restart = tkinter.Button(master, text="Начать заново",
                            command=prepare_and_start)
restart.pack()

prepare_and_start()
master.mainloop()


# --------------------------------------
WHITE = 1
BLACK = 2


                                  # Удобная функция для вычисления цвета противника
def opponent(color):
    if color == WHITE:
        return BLACK
    else:
        return WHITE


def print_board(board):  # Распечатать доску в текстовом виде (см. скриншот)
    print("     +----+----+----+----+----+----+----+----+")
    for row in range(7, -1, -1):
        print(" ", row, end="  ")
        for col in range(8):
            print("|", board.cell(row, col), end=" ")
        print("|")
        print("     +----+----+----+----+----+----+----+----+")
    print(end="        ")
    for col in range(8):
        print(col, end="    ")
    print()


def correct_coords(row, col):
    # внутри доски
    return 0 <= row < 8 and 0 <= col < 8


class Board:
    def __init__(self):
        self.color = WHITE
        self.field = []
        for row in range(8):
            self.field.append([None] * 8)
        self.field[0] = [
            Rook(WHITE),
            Knight(WHITE),
            Bishop(WHITE),
            Queen(WHITE),
            King(WHITE),
            Bishop(WHITE),
            Knight(WHITE),
            Rook(WHITE),
        ]
        self.field[1] = [
            Pawn(WHITE),
            Pawn(WHITE),
            Pawn(WHITE),
            Pawn(WHITE),
            Pawn(WHITE),
            Pawn(WHITE),
            Pawn(WHITE),
            Pawn(WHITE),
        ]
        self.field[6] = [
            Pawn(BLACK),
            Pawn(BLACK),
            Pawn(BLACK),
            Pawn(BLACK),
            Pawn(BLACK),
            Pawn(BLACK),
            Pawn(BLACK),
            Pawn(BLACK),
        ]
        self.field[7] = [
            Rook(BLACK),
            Knight(BLACK),
            Bishop(BLACK),
            Queen(BLACK),
            King(BLACK),
            Bishop(BLACK),
            Knight(BLACK),
            Rook(BLACK),
        ]

    def current_player_color(self):
        return self.color

    def cell(self, row, col):
        Where = self.field[row][col]
        if Where is None:
            return "  "
        color = Where.get_color()
        c = "w" if color == WHITE else "b"
        return c + Where.char()

    def get_Where(self, row, col):
        if correct_coords(row, col):
            return self.field[row][col]
        else:
            return None

    def move_Where(self, row, col, row1, col1):
        """Переместить фигуру из точки (row, col) в точку (row1, col1)"""

        if not correct_coords(row, col) or not correct_coords(row1, col1):
            return False
        if row == row1 and col == col1:
            return False  # нельзя пойти в ту же клетку
        Where = self.field[row][col]
        if Where is None:
            return False
        if Where.get_color() != self.color:
            return False
        if self.field[row1][col1] is None:
            if not Where.can_move(self, row, col, row1, col1):
                return False
        elif self.field[row1][col1].get_color() == opponent(Where.get_color()):
            if not Where.can_attack(self, row, col, row1, col1):
                return False
        else:
            return False
        self.field[row][col] = None  # Снять фигуру.
        self.field[row1][col1] = Where  # Поставить на новое место.
        self.color = opponent(self.color)
        return True

    def __str__(self):
        return str(print_board(self)).replace("None", "\r")


class Where:
    def __init__(self, color):
        self.color = color
        self.row = None
        self.col = None

    def can_move(self, board, row, col, row1, col1):
        if not correct_coords(row, col) or not correct_coords(row1, col1):
            return False
        if row == row1 and col == col1:
            return False
        if board.cell(row1, col1)[1] != " ":
            if board.get_Where(row1, col1).get_color() == self.color:
                return False
        return True

    def set_position(self, row, col):
        self.row = row
        self.col = col

    def get_color(self):
        return self.color

    def can_attack(self, board, row, col, row1, col1):
        return self.can_move(board, row, col, row1, col1)


class Rook(Where):
    def __init__(self, color):
        self.color = color

    def get_color(self):
        return self.color

    def char(self):
        return "R"

    def can_move(self, board, row, col, row1, col1):
        if row != row1 and col != col1:
            return False

        step = 1 if (row1 >= row) else -1
        for r in range(row + step, row1, step):    
            if not (board.get_Where(r, col) is None):
                return False

        step = 1 if (col1 >= col) else -1
        for c in range(col + step, col1, step):
            # Если на пути по вертикали есть фигура
            if not (board.get_Where(row, c) is None):
                return False

        return True

    def can_attack(self, board, row, col, row1, col1):
        return self.can_move(board, row, col, row1, col1)


class Pawn(Where):
    def __init__(self, color):
        self.color = color

    def get_color(self):
        return self.color

    def char(self):
        return "P"
    
    def can_move(self, board, row, col, row1, col1):
        if col != col1:
            return False
        if self.color == WHITE:
            direction = 1
            start_row = 1
        else:
            direction = -1
            start_row = 6

        # ход на 1 клетку
        if row + direction == row1:
            return True

        # ход на 2 клетки из начального положения
        if (
            row == start_row
            and row + 2 * direction == row1
            and board.field[row + direction][col] is None
        ):
            return True

        return False

    def can_attack(self, board, row, col, row1, col1):
        direction = 1 if (self.color == WHITE) else -1
        return row + direction == row1 and (col + 1 == col1 or col - 1 == col1)


class Knight(Where):
    def can_move(self, board, row, col, row1, col1):
        if not super().can_move(board, row, col, row1, col1):
            return False
        if not (
            (abs(row1 - row) == 2 and abs(col1 - col) == 1)
            or (abs(row1 - row) == 1 and abs(col1 - col) == 2)
        ):
            return False
        return True

    def char(self):
        return "N"


class King(Where):
    def can_move(self, board, row, col, row1, col1):
        if not super().can_move(board, row, col, row1, col1):
            return False
        if abs(row - row1) > 1 or abs(col - col1) > 1:
            return False
        return True

    def char(self):
        return "K"


class Queen(Where):
    def can_move(self, board, row, col, row1, col1):
        if not super().can_move(board, row, col, row1, col1):
            return False

        if abs(row1 - row) == abs(col1 - col):
            stepx = 1 if (col1 >= col) else -1
            stepy = 1 if (row1 >= row) else -1
            for i in range(1, abs(row1 - row)):
                if not (board.get_Where(row + i * stepy, col + i * stepx) is None):
                    return False
            return True
        if row == row1 or col == col1:
            step = 1 if (row1 >= row) else -1
            for r in range(row + step, row1, step):
                if not (board.get_Where(r, col) is None):
                    return False
            step = 1 if (col1 >= col) else -1
            for c in range(col + step, col1, step):
                if not (board.get_Where(row, c) is None):
                    return False
            return True
        return False

    def char(self):
        return "Q"


class Bishop(Where):
    def can_move(self, board, row, col, row1, col1):
        if not super().can_move(board, row, col, row1, col1):
            return False
        if not abs(row1 - row) == abs(col1 - col):
            return False

        stepx = 1 if (col1 >= col) else -1
        stepy = 1 if (row1 >= row) else -1

        for i in range(1, abs(row1 - row)):
            if not (board.get_Where(row + i * stepy, col + i * stepx) is None):
                return False
        return True

    def char(self):
        return "B"

board = Board()
print(board)